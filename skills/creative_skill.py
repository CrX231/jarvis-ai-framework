import os
import docx
import xlsxwriter
from pptx import Presentation
from core.skill_registry import BaseSkill

class CreativeSkill(BaseSkill):
    TRIGGERS = ["crea un documento", "monografía", "word", "crea un excel", "hoja de cálculo", "crea una presentación", "diapositivas", "powerpoint"]

    def __init__(self, context):
        super().__init__(context)
        self.brain = context.brain
        self.task_queue = context.task_queue
        self.event_bus = context.event_bus
        self.logger = context.logger
        
        self.output_folder = "documentos_generados"
        if not os.path.exists(self.output_folder):
            os.makedirs(self.output_folder)

    def execute(self, command, attachment_path=None):
        """Enruta al creador correcto según la palabra clave y lo lanza asíncronamente."""
        if any(word in command for word in ["excel", "hoja de cálculo"]):
            tema = command.replace("crea un excel sobre", "").replace("crea una hoja de cálculo sobre", "").replace("excel", "").strip()
            self.task_queue.add_task(self._create_excel_async, tema)
            return "Generando hoja de cálculo en segundo plano."
            
        elif any(word in command for word in ["presentación", "diapositivas", "powerpoint"]):
            tema = command.replace("crea una presentación sobre", "").replace("crea unas diapositivas sobre", "").replace("powerpoint", "").strip()
            self.task_queue.add_task(self._create_pptx_async, tema)
            return "Diseñando diapositivas en segundo plano."
            
        else: # Por defecto es Word/Monografía
            tema = command.replace("crea un documento de word sobre", "").replace("crea un documento sobre", "").replace("haz una monografía sobre", "").replace("word", "").strip()
            self.task_queue.add_task(self._create_word_async, tema)
            return "Iniciando redacción del documento en segundo plano. Le notificaré cuando esté listo."

    def _ask_gemini_deep(self, prompt):
        try:
            response = self.brain.client.models.generate_content(
                model=self.brain.model_id,
                contents=prompt
            )
            return response.text.replace("*", "")
        except Exception as e:
            self.logger.error(f"Error de IA: {e}")
            return ""

    def _create_word_async(self, topic):
        self.logger.info(f"[Creative] Redactando monografía sobre: {topic}")
        prompt = f"Escribe una monografía extensa y profesional sobre '{topic}'. Estructura con 'H1: Título' y 'H2: Subtítulo'. No uses markdown."
        contenido = self._ask_gemini_deep(prompt)
        if not contenido: return
        
        doc = docx.Document()
        doc.add_heading(topic.upper(), level=0)
        
        for linea in contenido.split('\n'):
            linea = linea.strip()
            if linea.startswith("H1:"):
                doc.add_heading(linea.replace("H1:", ""), level=1)
            elif linea.startswith("H2:"):
                doc.add_heading(linea.replace("H2:", ""), level=2)
            elif linea:
                doc.add_paragraph(linea)
        
        filename = f"Monografia_{topic.replace(' ', '_')[:20]}.docx"
        path = os.path.join(self.output_folder, filename)
        doc.save(path)
        os.startfile(path)
        self.event_bus.publish("SPEAK_REQUEST", {"text": "Documento completado y abierto en su pantalla, señor."})

    def _create_excel_async(self, topic):
        self.logger.info(f"[Creative] Generando hoja de cálculo sobre: {topic}")
        prompt = f"Genera una lista de datos para un Excel sobre '{topic}'. Devuelve los datos en formato CSV (separados por comas), máximo 10 filas. Solo los datos, nada de texto extra."
        datos_raw = self._ask_gemini_deep(prompt)
        if not datos_raw: return
        
        filename = f"Excel_{topic.replace(' ', '_')[:20]}.xlsx"
        path = os.path.join(self.output_folder, filename)
        
        workbook = xlsxwriter.Workbook(path)
        worksheet = workbook.add_worksheet()
        bold = workbook.add_format({'bold': True, 'bg_color': '#D4AF37', 'font_color': 'white'})
        
        for r, linea in enumerate(datos_raw.split('\n')):
            if not linea.strip(): continue
            columnas = linea.split(',')
            for c, valor in enumerate(columnas):
                if r == 0:
                    worksheet.write(r, c, valor.strip(), bold)
                else:
                    worksheet.write(r, c, valor.strip())
        
        workbook.close()
        os.startfile(path)
        self.event_bus.publish("SPEAK_REQUEST", {"text": "Hoja de cálculo generada y lista en su pantalla."})

    def _create_pptx_async(self, topic):
        self.logger.info(f"[Creative] Diseñando presentación sobre: {topic}")
        prompt = f"Genera el esquema para una presentación de 5 diapositivas sobre '{topic}'. Para cada diapositiva escribe 'T: Titulo' y luego 'C: Contenido'. No uses markdown."
        esquema = self._ask_gemini_deep(prompt)
        if not esquema: return
        
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = topic.upper()
        slide.placeholders[1].text = "Generado por Jarvis AI"

        current_slide = None
        for linea in esquema.split('\n'):
            linea = linea.strip()
            if linea.startswith("T:"):
                bullet_slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                current_slide.shapes.title.text = linea.replace("T:", "").strip()
            elif linea.startswith("C:") and current_slide:
                current_slide.placeholders[1].text += linea.replace("C:", "").strip() + "\n"

        filename = f"Presentacion_{topic.replace(' ', '_')[:20]}.pptx"
        path = os.path.join(self.output_folder, filename)
        prs.save(path)
        os.startfile(path)
        self.event_bus.publish("SPEAK_REQUEST", {"text": "Las diapositivas han sido diseñadas y abiertas, señor."})